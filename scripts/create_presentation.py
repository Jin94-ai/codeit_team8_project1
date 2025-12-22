"""
Health Eat 발표자료 생성 스크립트
20장 PPTX 자동 생성
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# 색상 정의
NAVY = RGBColor(0x1a, 0x36, 0x5d)  # 네이비
WHITE = RGBColor(0xff, 0xff, 0xff)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xf5, 0xf5, 0xf5)
ACCENT = RGBColor(0x00, 0x7a, 0xcc)  # 파란색 강조

def add_title_slide(prs, title, subtitle=""):
    """표지 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 빈 슬라이드

    # 배경색
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = NAVY
    background.line.fill.background()

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # 부제목
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, content_lines, is_bullet=True):
    """일반 컨텐츠 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 상단 바
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = NAVY
    header.line.fill.background()

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # 내용
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, line in enumerate(content_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if is_bullet and line.strip():
            p.text = line
            p.level = 0
        else:
            p.text = line

        p.font.size = Pt(20)
        p.font.color.rgb = GRAY
        p.space_after = Pt(12)

    return slide

def add_table_slide(prs, title, headers, rows):
    """테이블 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 상단 바
    header_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = NAVY
    header_shape.line.fill.background()

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # 테이블
    num_cols = len(headers)
    num_rows = len(rows) + 1

    table = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(0.5), Inches(1.6),
        Inches(9), Inches(0.5 * num_rows)
    ).table

    # 헤더
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.font.size = Pt(16)
        p.alignment = PP_ALIGN.CENTER

    # 데이터
    for row_idx, row in enumerate(rows):
        for col_idx, cell_text in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(14)
            p.font.color.rgb = GRAY
            p.alignment = PP_ALIGN.CENTER

    return slide

def add_architecture_slide(prs, title):
    """아키텍처 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 상단 바
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = NAVY
    header.line.fill.background()

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Stage 1 박스
    box1 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.8), Inches(4), Inches(2)
    )
    box1.fill.solid()
    box1.fill.fore_color.rgb = RGBColor(0xe3, 0xf2, 0xfd)

    text1 = slide.shapes.add_textbox(Inches(0.7), Inches(1.9), Inches(3.6), Inches(1.8))
    tf = text1.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Stage 1: YOLO11m Detector"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = NAVY

    for line in ["Input: 원본 이미지", "Output: Bounding Boxes", "클래스: 단일 (Pill)"]:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.font.color.rgb = GRAY

    # 화살표
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, Inches(4.6), Inches(2.5), Inches(0.8), Inches(0.5)
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = NAVY

    # Stage 2 박스
    box2 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.5), Inches(1.8), Inches(4), Inches(2)
    )
    box2.fill.solid()
    box2.fill.fore_color.rgb = RGBColor(0xe8, 0xf5, 0xe9)

    text2 = slide.shapes.add_textbox(Inches(5.7), Inches(1.9), Inches(3.6), Inches(1.8))
    tf = text2.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Stage 2: ConvNeXt Classifier"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = NAVY

    for line in ["Input: 크롭 이미지 (224x224)", "Output: K-code + Confidence", "클래스: 74개"]:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.font.color.rgb = GRAY

    # 하단 설명
    desc = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1.5))
    tf = desc.text_frame
    tf.word_wrap = True

    for line in ["Detection과 Classification 분리로 각각 최적화",
                 "Detector: 위치 검출만 담당 (일반화 용이)",
                 "Classifier: 분류만 담당 (독립 최적화)"]:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
        p.font.color.rgb = GRAY
        p.space_before = Pt(8)

    return slide

def add_score_history_slide(prs):
    """점수 히스토리 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 상단 바
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = NAVY
    header.line.fill.background()

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "실험 결과"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # 테이블
    headers = ["#", "Score", "설명", "변화"]
    rows = [
        ["1", "0.82", "Baseline (YOLO12m)", "-"],
        ["-", "0.6", "복합경구제 (bbox 1개)", "-0.22"],
        ["2", "0.822", "단일경구제", "+0.222"],
        ["3", "0.920", "2-Stage 도입", "+0.098"],
        ["4", "0.963", "데이터셋 개선", "+0.043"],
        ["5", "0.96703", "데이터셋 정제 (3~4개 bbox)", "+0.004"],
    ]

    table = slide.shapes.add_table(
        len(rows) + 1, len(headers),
        Inches(0.3), Inches(1.5),
        Inches(9.4), Inches(3.5)
    ).table

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.font.size = Pt(14)
        p.alignment = PP_ALIGN.CENTER

    for row_idx, row in enumerate(rows):
        for col_idx, val in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.alignment = PP_ALIGN.CENTER

            # 최종 점수 강조
            if row_idx == len(rows) - 1:
                p.font.bold = True
                p.font.color.rgb = ACCENT

    return slide

def add_conclusion_slide(prs):
    """결론 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경색
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = NAVY
    background.line.fill.background()

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "결론"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # 핵심 메시지
    result_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1))
    tf = result_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Baseline 0.82 → Best 0.96703 (+18%)"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x4c, 0xaf, 0x50)  # 녹색
    p.alignment = PP_ALIGN.CENTER

    # 핵심 포인트
    points = [
        "End-to-end의 한계 인식 → 2-Stage 전환",
        "수많은 시행착오가 성공의 밑거름",
        "데이터 품질 > 데이터 양",
        "비즈니스 방향성 고민 → B2B SaaS 모델 검토"
    ]

    content_box = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(8), Inches(2.5))
    tf = content_box.text_frame

    for i, point in enumerate(points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{i+1}. {point}"
        p.font.size = Pt(20)
        p.font.color.rgb = WHITE
        p.space_after = Pt(16)

    return slide

def add_qa_slide(prs):
    """Q&A 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경색
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = NAVY
    background.line.fill.background()

    # 감사합니다
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "감사합니다"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # GitHub
    github_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(0.5))
    tf = github_box.text_frame
    p = tf.paragraphs[0]
    p.text = "github.com/Jin94-ai/codeit_team8_project1"
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    return slide

def create_presentation():
    """20장 발표자료 생성"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 1. 표지
    add_title_slide(prs,
        "Health Eat - AI 알약 인식",
        "코드잇 8팀 | 2025.12.04 ~ 12.23"
    )

    # 2. 목차
    add_content_slide(prs, "목차", [
        "1. 프로젝트 개요",
        "2. 데이터 분석",
        "3. 시행착오의 여정",
        "4. 핵심 전환점: 2-Stage Pipeline",
        "5. 최종 모델 아키텍처",
        "6. 실험 결과",
        "7. 팀 협업 & 비즈니스 방향성",
        "8. 결론 및 회고"
    ])

    # 3. 프로젝트 개요
    add_table_slide(prs, "프로젝트 개요",
        ["항목", "내용"],
        [
            ["목표", "Kaggle mAP@[0.75:0.95] 최대화"],
            ["기간", "3주 (12/04 ~ 12/23)"],
            ["최종 점수", "0.96703"],
            ["Baseline 대비", "+18% 개선"]
        ]
    )

    # 4. 데이터 분석
    add_table_slide(prs, "데이터 분석 - 문제 인식",
        ["구분", "학습 데이터", "테스트 데이터"],
        [
            ["이미지 수", "232개", "843개"],
            ["클래스 수", "56개", "미공개"],
            ["문제점", "클래스 불균형 (1:80)", ""]
        ]
    )

    # 5. AIHub 데이터 시도
    add_table_slide(prs, "데이터 확장 시도 - AIHub",
        ["시도", "데이터셋", "결과"],
        [
            ["1차", "복합경구제", "실패 (0.6)"],
            ["2차", "단일경구제", "미미 (0.822)"],
            ["3차", "18개 추가 클래스", "74개 클래스 확보"]
        ]
    )

    # 6. 시행착오 - 복합경구제
    add_content_slide(prs, "시행착오: AIHub 복합경구제 실패", [
        "문제: 56개 클래스 필터링",
        "  - 이미지당 라벨 1개만 추출",
        "  - bbox 1개만 검출 → Score 0.6",
        "",
        "원인:",
        "  - 데이터 구조 미이해",
        "  - 필터링 로직 오류",
        "",
        "교훈: 데이터 구조 이해 필수"
    ], is_bullet=False)

    # 7. 핵심 전환점
    add_content_slide(prs, "핵심 전환점: 2-Stage Pipeline", [
        "기존 방식의 한계:",
        "  이미지 → YOLO → 74개 클래스 직접 예측 (End-to-end)",
        "  → 클래스 불균형, 데이터 부족으로 성능 한계",
        "",
        "새로운 접근:",
        "  이미지 → Detector → 알약 위치 검출 (단일 클래스)",
        "         → Classifier → 74개 클래스 분류",
        "",
        "분리의 이점:",
        "  - Detection과 Classification 독립 최적화",
        "  - 각각 최적의 모델 선택 가능"
    ], is_bullet=False)

    # 8. 아키텍처
    add_architecture_slide(prs, "최종 모델 아키텍처: 2-Stage Pipeline")

    # 9. Stage 1 상세
    add_table_slide(prs, "Stage 1: YOLO11m Detector",
        ["설정", "값"],
        [
            ["모델", "yolo11m.pt"],
            ["imgsz", "640"],
            ["mAP50", "0.995"],
            ["mAP50-95", "0.85"],
            ["Precision", "0.99"],
            ["Recall", "0.99"]
        ]
    )

    # 10. Stage 2 상세
    add_table_slide(prs, "Stage 2: ConvNeXt Classifier",
        ["설정", "값"],
        [
            ["모델", "convnext_tiny (ImageNet)"],
            ["img_size", "224"],
            ["Val Accuracy", "98.5%"],
            ["Early Stop", "Epoch 21"],
            ["비고", "YOLO-cls와 분류 정확도 비슷"]
        ]
    )

    # 11. 데이터 정제
    add_table_slide(prs, "데이터 정제 - Best Score 달성",
        ["정제 기준", "값", "이유"],
        [
            ["bbox 개수", "2~4개", "테스트와 동일"],
            ["IoU 임계값", "0.7", "진짜 중복만"],
            ["bbox 크기", "30px 이상", "너무 작은 것 제외"],
            ["종횡비", "3.5 이하", "비정상 형태 제외"]
        ]
    )

    # 12. 실험 결과
    add_score_history_slide(prs)

    # 13. 실패 사례
    add_table_slide(prs, "실패 사례와 교훈",
        ["시도", "결과", "교훈"],
        [
            ["복합경구제", "bbox 1개 검출", "데이터 구조 이해 필수"],
            ["ROI Crop", "검출 불가", "크롭 로직 검증 필요"],
            ["imgsz 1280", "0.713", "학습/추론 설정 일치"],
            ["TTA 적용", "0.533", "bbox 정밀도 저하"]
        ]
    )

    # 14. 팀 협업
    add_table_slide(prs, "팀 협업 & 역할 분담",
        ["담당자", "역할", "주요 기여"],
        [
            ["이진석", "Leader", "2-Stage Pipeline, 데이터셋 정제"],
            ["김민우", "Data Engineer", "서비스 타당성 분석"],
            ["김나연", "Data Engineer", "B2B 시장 조사"],
            ["김보윤", "Model Architect", "Cloud 아키텍처"],
            ["황유민", "Experiment Lead", "챗봇 프롬프트"]
        ]
    )

    # 15. 비즈니스 방향성
    add_table_slide(prs, "비즈니스 방향성",
        ["시장", "타겟", "가치 제안"],
        [
            ["B2C", "고령자", "알약 식별, 효능 안내"],
            ["B2B", "응급구조대", "신속한 약물 식별"],
            ["B2B", "수사기관", "범죄현장 약물 식별"]
        ]
    )

    # 16. 서비스 아키텍처
    add_table_slide(prs, "서비스 아키텍처: Cloud 방식",
        ["항목", "On-Device", "Cloud (채택)"],
        [
            ["앱 크기", "~100MB", "~20MB"],
            ["모델 업데이트", "앱 재배포", "서버만 수정"],
            ["기기 호환성", "고사양 필요", "저사양 가능"]
        ]
    )

    # 17. 타임라인
    add_content_slide(prs, "프로젝트 타임라인", [
        "Phase 1-2 (12/04~09): 셋업, EDA, 56개 클래스 확인",
        "Phase 3 (12/09~12): Baseline 0.82 제출",
        "Phase 4 (12/11~15): AIHub 시도 - 복합경구제 실패, 단일경구제",
        "Phase 5 (12/17~18): 18개 추가 클래스 → 74개",
        "Phase 6 (12/18): 2-Stage Pipeline → 0.920 → 0.963",
        "Phase 7 (12/18~19): 데이터 정제 → 0.96703",
        "Phase 8 (12/19~22): 비즈니스 논의, 마무리"
    ])

    # 18. 핵심 성공 요인
    add_content_slide(prs, "핵심 성공 요인", [
        "1. 2-Stage 분리 (+0.10)",
        "   Detection과 Classification 독립 최적화",
        "",
        "2. 데이터셋 개선 (+0.043)",
        "   AIHub 데이터 추가, 테스트 환경과 유사하게 조정",
        "",
        "3. 데이터셋 정제 (+0.004)",
        "   테스트와 유사한 3~4개 bbox 이미지만 사용",
        "",
        "4. 빠른 실험 사이클",
        "   실패 → 원인 분석 → 방향 전환"
    ], is_bullet=False)

    # 19. 결론
    add_conclusion_slide(prs)

    # 20. Q&A
    add_qa_slide(prs)

    # 저장
    output_path = os.path.join(os.path.dirname(__file__), "..", "docs", "Health_Eat_Presentation.pptx")
    prs.save(output_path)
    print(f"발표자료 생성 완료: {output_path}")
    return output_path

if __name__ == "__main__":
    create_presentation()
